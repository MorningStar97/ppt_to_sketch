import os
import json
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import uuid
import logging
import io
import base64

logger = logging.getLogger(__name__)

class PPTToSketchConverter:
    """
    PPT转Sketch转换器 - 简化版，无缩放逻辑
    """
    
    def __init__(self, verbose=False):
        """
        初始化转换器
        """
        self.verbose = verbose
        self.image_dict = {}
        self.artboard_width = None
        self.artboard_height = None
        
    def log(self, message, level='info'):
        """日志记录"""
        if self.verbose:
            if level == 'error':
                logger.error(message)
            elif level == 'warning':
                logger.warning(message)
            else:
                logger.info(message)
    
    def extract_color(self, color_obj):
        """
        提取颜色信息并转换为 Sketch 格式
        
        Args:
            color_obj: PPT 颜色对象
            
        Returns:
            dict: Sketch 颜色格式 (RGBA 0-1 范围)
        """
        try:
            if isinstance(color_obj, RGBColor):
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": color_obj[2] / 255.0,
                    "green": color_obj[1] / 255.0,
                    "red": color_obj[0] / 255.0
                }
            elif hasattr(color_obj, 'rgb') and color_obj.rgb:
                rgb = color_obj.rgb
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": rgb[2] / 255.0,
                    "green": rgb[1] / 255.0,
                    "red": rgb[0] / 255.0
                }
            elif isinstance(color_obj, (tuple, list)) and len(color_obj) >= 3:
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": color_obj[2] / 255.0,
                    "green": color_obj[1] / 255.0,
                    "red": color_obj[0] / 255.0
                }
            elif isinstance(color_obj, int):
                # 处理整数颜色值
                red = (color_obj >> 16) & 0xFF
                green = (color_obj >> 8) & 0xFF
                blue = color_obj & 0xFF
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": blue / 255.0,
                    "green": green / 255.0,
                    "red": red / 255.0
                }
        except Exception as e:
            self.log(f"颜色提取失败: {str(e)}", 'warning')
        
        # 默认返回黑色
        return {
            "_class": "color",
            "alpha": 1,
            "blue": 0,
            "green": 0,
            "red": 0
        }

    def optimize_image(self, image_blob, max_size=4096):
        """优化图片大小，转换为JPEG格式，并限制最大尺寸"""
        try:
            image = Image.open(io.BytesIO(image_blob))

            # 统一转换为RGB，避免RGBA带来的透明度问题
            if image.mode == 'RGBA':
                background = Image.new('RGB', image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[-1] if len(image.split()) == 4 else None)
                image = background
            elif image.mode != 'RGB':
                image = image.convert('RGB')

            # 如果图片尺寸过大，进行缩放
            if max(image.size) > max_size:
                image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
                self.log(f"图片尺寸已优化至 {image.size}")

            # 保存为高质量JPEG
            output = io.BytesIO()
            image.save(output, format='JPEG', quality=90, optimize=True)
            optimized_data = output.getvalue()
            
            self.log(f"图片优化: {len(image_blob)} -> {len(optimized_data)} 字节 (JPEG)")
            return optimized_data
            
        except Exception as e:
            self.log(f"图片优化失败: {str(e)}", 'warning')
            return image_blob

    def create_text_layer(self, shape, layer_name):
        """创建文本图层 - 直接坐标映射，让artboard处理裁剪"""
        try:
            # 直接使用点坐标，无任何缩放
            left = shape.left.pt if shape.left else 0
            top = shape.top.pt if shape.top else 0
            width = shape.width.pt if shape.width else 100
            height = shape.height.pt if shape.height else 50
            
            # 修正旋转角度：PPT和Sketch的坐标系方向相反
            rotation = -shape.rotation if hasattr(shape, 'rotation') else 0
            
            text_content = ""
            font_name = "Arial"
            font_size = 12 # 默认值
            text_color = self.extract_color((0, 0, 0))
            alignment = 0 # 左对齐

            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                text_content = text_frame.text

                # 智能字体样式提取：找到应用范围最广的样式
                style_counts = {}
                for p in text_frame.paragraphs:
                    # 优先使用段落级别的对齐方式
                    align_map = {PP_ALIGN.LEFT: 0, PP_ALIGN.CENTER: 2, PP_ALIGN.RIGHT: 1, PP_ALIGN.JUSTIFY: 3}
                    alignment = align_map.get(p.alignment, alignment)

                    for run in p.runs:
                        if run.font.size and run.font.name:
                            # 直接使用字体点大小，无缩放
                            size = run.font.size.pt or 12
                            name = run.font.name
                            color = self.extract_color(run.font.color.rgb) if run.font.color and hasattr(run.font.color, 'rgb') else text_color
                            
                            style_key = (size, name, json.dumps(color))
                            style_counts[style_key] = style_counts.get(style_key, 0) + len(run.text)
                
                # 选择字符数最多的样式作为主样式
                if style_counts:
                    dominant_style = max(style_counts, key=style_counts.get)
                    font_size = dominant_style[0]
                    font_name = dominant_style[1]
                    text_color = json.loads(dominant_style[2])
            
            layer = {
                "_class": "text",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {"_class": "rect", "constrainProportions": False, "height": height, "width": width, "x": left, "y": top},
                "rotation": rotation,
                "isVisible": True,
                "isLocked": False,
                "attributedString": {
                    "_class": "attributedString",
                    "string": text_content,
                    "attributes": [{
                        "_class": "stringAttribute",
                        "location": 0, "length": len(text_content),
                        "attributes": {
                            "MSAttributedStringFontAttribute": {"_class": "fontDescriptor", "attributes": {"name": font_name, "size": font_size}},
                            "MSAttributedStringColorAttribute": text_color,
                            "paragraphStyle": {"_class": "paragraphStyle", "alignment": alignment}
                        }
                    }]
                },
                "style": {"_class": "style", "endDecorationType": 0, "miterLimit": 10, "startDecorationType": 0, "windingRule": 1}
            }
            return layer
        except Exception as e:
            self.log(f"创建文本图层失败: {layer_name} - {e}", 'error')
            return None

    def create_image_layer(self, shape, layer_name):
        """创建图片图层 - 直接坐标映射，让artboard处理裁剪"""
        try:
            if not hasattr(shape, 'image'):
                return None

            # 直接使用点坐标，无任何缩放
            left = shape.left.pt if shape.left else 0
            top = shape.top.pt if shape.top else 0
            width = shape.width.pt if shape.width else 100
            height = shape.height.pt if shape.height else 100

            # 修正旋转角度：PPT和Sketch的坐标系方向相反
            rotation = -shape.rotation if hasattr(shape, 'rotation') else 0

            # 调试输出旋转信息
            if self.verbose and rotation != 0:
                self.log(f"🔄 图片旋转信息 - 名称: {layer_name}, PPT角度: {shape.rotation}°, Sketch角度: {rotation}°")

            # 保存图片数据 - 使用用户验证过的格式
            image_data = shape.image.blob
            image_ref = f"images/{str(uuid.uuid4())}.png"
            self.image_dict[image_ref] = {"type": "Buffer", "data": list(image_data)}
            
            layer = {
                "_class": "bitmap",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {"_class": "rect", "constrainProportions": True, "height": height, "width": width, "x": left, "y": top},
                "rotation": rotation,  # 修正后的角度值
                "isVisible": True,
                "isLocked": False,
                "image": {"_class": "MSJSONFileReference", "_ref_class": "MSImageData", "_ref": image_ref},
                "style": {"_class": "style", "endDecorationType": 0, "miterLimit": 10, "startDecorationType": 0, "windingRule": 1}
            }
            return layer
        except Exception as e:
            self.log(f"创建图片图层失败: {layer_name} - {e}", 'error')
            return None

    def create_shape_layer(self, shape, layer_name):
        """创建形状图层 - 直接坐标映射，让artboard处理裁剪"""
        try:
            # 直接使用点坐标，无任何缩放
            left = shape.left.pt if shape.left else 0
            top = shape.top.pt if shape.top else 0
            width = shape.width.pt if shape.width else 100
            height = shape.height.pt if shape.height else 100

            # 修正旋转角度：PPT和Sketch的坐标系方向相反
            rotation = -shape.rotation if hasattr(shape, 'rotation') else 0

            fill_color = self.extract_color((128, 128, 128)) # 默认灰色
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                fill_color = self.extract_color(shape.fill.fore_color)
            
            path_points = [
                {"_class": "curvePoint", "cornerRadius": 0, "curveFrom": "{0, 0}", "curveTo": "{0, 0}", "hasCurveFrom": False, "hasCurveTo": False, "point": "{0, 0}"},
                {"_class": "curvePoint", "cornerRadius": 0, "curveFrom": "{1, 0}", "curveTo": "{1, 0}", "hasCurveFrom": False, "hasCurveTo": False, "point": "{1, 0}"},
                {"_class": "curvePoint", "cornerRadius": 0, "curveFrom": "{1, 1}", "curveTo": "{1, 1}", "hasCurveFrom": False, "hasCurveTo": False, "point": "{1, 1}"},
                {"_class": "curvePoint", "cornerRadius": 0, "curveFrom": "{0, 1}", "curveTo": "{0, 1}", "hasCurveFrom": False, "hasCurveTo": False, "point": "{0, 1}"}
            ]

            layer = {
                "_class": "rectangle",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {"_class": "rect", "constrainProportions": False, "height": height, "width": width, "x": left, "y": top},
                "rotation": rotation,  # 修正后的角度值
                "isVisible": True,
                "isLocked": False,
                "path": {"_class": "path", "isClosed": True, "pointRadiusBehaviour": 1, "points": path_points},
                "style": {
                    "_class": "style", "endDecorationType": 0, "miterLimit": 10, "startDecorationType": 0, "windingRule": 1,
                    "fills": [{"_class": "fill", "isEnabled": True, "color": fill_color, "fillType": 0}]
                }
            }
            return layer
        except Exception as e:
            self.log(f"创建形状图层失败: {layer_name} - {e}", 'error')
            return None

    def create_group_layer(self, shape, layer_name):
        """创建组图层 - 修正相对坐标计算"""
        try:
            # 直接使用点坐标，无任何缩放
            group_left = shape.left.pt if shape.left else 0
            group_top = shape.top.pt if shape.top else 0
            width = shape.width.pt if shape.width else 100
            height = shape.height.pt if shape.height else 100

            # 修正旋转角度：PPT和Sketch的坐标系方向相反
            rotation = -shape.rotation if hasattr(shape, 'rotation') else 0
            
            sub_layers = []
            if hasattr(shape, 'shapes'):
                # 严格按照PPT内部顺序处理子图层
                for i, sub_shape in enumerate(shape.shapes):
                    sub_layer = self.process_shape(sub_shape, f"{layer_name}_child_{i}")
                    if sub_layer:
                        # 正确计算相对于父组的坐标
                        if 'frame' in sub_layer:
                            sub_layer['frame']['x'] -= group_left
                            sub_layer['frame']['y'] -= group_top
                        sub_layers.append(sub_layer)
            
            if not sub_layers:
                self.log(f"跳过空组: {layer_name}", 'warning')
                return None

            layer = {
                "_class": "group",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {"_class": "rect", "constrainProportions": False, "height": height, "width": width, "x": group_left, "y": group_top},
                "rotation": rotation,  # 修正后的角度值
                "isVisible": True,
                "isLocked": False,
                "layers": sub_layers,
                "style": {"_class": "style", "endDecorationType": 0, "miterLimit": 10, "startDecorationType": 0, "windingRule": 1}
            }
            return layer
        except Exception as e:
            self.log(f"创建组图层失败: {layer_name} - {e}", 'error')
            return None

    def extract_slide_background(self, slide):
        """提取幻灯片背景信息"""
        try:
            bg_color = None
            if hasattr(slide, 'background') and hasattr(slide.background, 'fill') and hasattr(slide.background.fill, 'fore_color'):
                bg_color = self.extract_color(slide.background.fill.fore_color)
            
            if bg_color:
                return {
                    "_class": "rectangle",
                    "do_objectID": str(uuid.uuid4()),
                    "name": "Slide Background",
                    "frame": {"_class": "rect", "height": self.artboard_height, "width": self.artboard_width, "x": 0, "y": 0},
                    "isLocked": True,
                    "style": {"_class": "style", "fills": [{"_class": "fill", "isEnabled": True, "color": bg_color, "fillType": 0}]}
                }
        except Exception:
            pass # 背景不是纯色，忽略
        return None

    def process_shape(self, shape, layer_name):
        """处理单个形状，根据类型分派 - 保真版"""
        try:
            if not shape.width or not shape.height:
                self.log(f"跳过无尺寸的形状: {layer_name}", 'warning')
                return None

            shape_type = shape.shape_type
            
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                return self.create_group_layer(shape, layer_name)
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self.create_image_layer(shape, layer_name)
            elif hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                return self.create_text_layer(shape, layer_name)
            else: # 其他所有类型都视为基本形状
                return self.create_shape_layer(shape, layer_name)
                
        except Exception as e:
            self.log(f"处理形状失败: {layer_name} - {e}", 'error')
            return None
    
    def convert_slide_to_artboard(self, slide, slide_index):
        """将PPT幻灯片转换为Sketch画板 - 层级保真"""
        try:
            artboard_name = f"Slide {slide_index + 1}"
            layers = []
            
            # 1. 添加幻灯片背景色（如果存在）
            slide_bg_layer = self.extract_slide_background(slide)
            if slide_bg_layer:
                layers.append(slide_bg_layer)
            
            # 2. 严格按照z-order遍历所有形状
            # python-pptx的slide.shapes本身就是从底层到顶层的顺序
            for i, shape in enumerate(slide.shapes):
                layer = self.process_shape(shape, f"Layer_{i}")
                if layer:
                    layers.append(layer)
            
            # 创建画板 - 补全所有标准属性
            artboard = {
                "_class": "artboard",
                "do_objectID": str(uuid.uuid4()),
                "name": artboard_name,
                "frame": {
                    "_class": "rect",
                    "constrainProportions": False,
                    "height": self.artboard_height,
                    "width": self.artboard_width,
                    "x": 0,
                    "y": slide_index * (self.artboard_height + 100)
                },
                "layers": layers,
                "hasBackgroundColor": True,
                "backgroundColor": self.extract_color((255, 255, 255)),
                "hasClippingMask": True,
                "clippingMaskMode": 0,
                # -- 补全标准属性 --
                "booleanOperation": -1,
                "isFixedToViewport": False,
                "isFlippedHorizontal": False,
                "isFlippedVertical": False,
                "isLocked": False,
                "isVisible": True,
                "layerListExpandedType": 0,
                "nameIsFixed": False,
                "resizingConstraint": 63,
                "resizingType": 0,
                "rotation": 0,
                "shouldBreakMaskChain": False,
                "exportOptions": {
                    "_class": "exportOptions",
                    "includedLayerIds": [],
                    "layerOptions": 0,
                    "shouldTrim": False,
                    "exportFormats": []
                },
                "style": {
                    "_class": "style",
                    "endMarkerType": 0,
                    "miterLimit": 10,
                    "startMarkerType": 0,
                    "windingRule": 1,
                    "borders": [],
                    "fills": [],
                    "shadows": []
                }
            }
            return artboard
        except Exception as e:
            self.log(f"转换幻灯片失败: {slide_index} - {e}", 'error')
            return None

    def convert_ppt_to_sketch(self, ppt_file_path, output_dir):
        """将PPT文件转换为Sketch格式 - 纯点单位版"""
        try:
            self.log(f"开始转换: {ppt_file_path}")
            presentation = Presentation(ppt_file_path)
            
            # 直接使用点单位属性，避免EMU转换
            self.artboard_width = presentation.slide_width.pt
            self.artboard_height = presentation.slide_height.pt
            self.log(f"画板尺寸: {self.artboard_width:.2f} x {self.artboard_height:.2f} 点")

            artboards = []
            for i, slide in enumerate(presentation.slides):
                artboard = self.convert_slide_to_artboard(slide, i)
                if artboard:
                    artboards.append(artboard)
            
            if not artboards:
                raise Exception("未能成功转换任何幻灯片")
            
            sketch_data = self._create_sketch_document(artboards)
            return self._generate_sketch_file(sketch_data, output_dir)
            
        except Exception as e:
            self.log(f"转换过程发生严重错误: {e}", 'error')
            raise e
    
    def _create_sketch_document(self, artboards):
        """
        创建Sketch文档结构 - 基于标准文件进行精确重写
        """
        
        # 1. 创建页面数据和页面引用
        pages_data = []
        page_refs = []
        
        for artboard in artboards:
            page_id = str(uuid.uuid4())
            page_data = {
                "_class": "page",
                "do_objectID": page_id,
                "booleanOperation": -1, "isFixedToViewport": False, "isFlippedHorizontal": False, "isFlippedVertical": False,
                "isLocked": False, "isVisible": True, "layerListExpandedType": 0,
                "name": artboard.get("name", "Page 1"), "nameIsFixed": False,
                "resizingConstraint": 63, "resizingType": 0, "rotation": 0, "shouldBreakMaskChain": False,
                "exportOptions": {"_class": "exportOptions", "includedLayerIds": [], "layerOptions": 0, "shouldTrim": False, "exportFormats": []},
                "frame": {"_class": "rect", "constrainProportions": False, "height": self.artboard_height, "width": self.artboard_width, "x": 0, "y": 0},
                "clippingMaskMode": 0, "hasClippingMask": False,
                "style": {"_class": "style", "endMarkerType": 0, "miterLimit": 10, "startMarkerType": 0, "windingRule": 1, "borders": [], "fills": [], "shadows": []},
                "layers": [artboard]
            }
            pages_data.append(page_data)
            page_refs.append({"_class": "MSJSONFileReference", "_ref_class": "MSImmutablePage", "_ref": f"pages/{page_id}"})

        # 2. 创建 document.json
        document = {
            "_class": "document",
            "do_objectID": str(uuid.uuid4()),
            "assets": {"_class": "assetCollection", "images": [], "colorAssets": [], "exportPresets": [], "gradientAssets": [], "colors": [], "gradients": []},
            "colorSpace": 1,
            "currentPageIndex": 0,
            "foreignLayerStyles": [], "foreignSymbols": [], "foreignTextStyles": [], "foreignSwatches": [],
            "layerStyles": {"_class": "sharedStyleContainer", "objects": []},
            "layerTextStyles": {"_class": "sharedTextStyleContainer", "objects": []},
            "pages": page_refs,
            "perDocumentLibraries": [],
            "sharedSwatches": {"_class": "swatchContainer", "objects": []}
        }

        # 3. 创建 meta.json
        meta = {
            "commit": "local",
            "pagesAndArtboards": {},
            "version": 168,
            "compatibilityVersion": 99,
            "app": "com.bohemiancoding.sketch3",
            "autosaved": 0, "variant": "NONAPPSTORE",
            "created": {"appVersion": "2025.1.3", "build": 203472, "app": "com.bohemiancoding.sketch3", "commit": "local", "version": 168, "variant": "NONAPPSTORE"},
            "saveHistory": ["NONAPPSTORE.203472"],
            "appVersion": "2025.1.3", "build": 203472
        }
        for page in pages_data:
            artboard_info = {}
            for layer in page.get("layers", []):
                if layer.get("_class") == "artboard":
                    artboard_info[layer["do_objectID"]] = {"name": layer["name"]}
            meta["pagesAndArtboards"][page["do_objectID"]] = {"name": page["name"], "artboards": artboard_info}

        # 4. 创建 user.json
        user = {page["do_objectID"]: {"scrollOrigin": "{0, 0}", "zoomValue": 1} for page in pages_data}

        return {
            "contents": {
                "document": document,
                "pages": pages_data,
                "meta": meta,
                "user": user,
                "workspace": {}
            },
            "imageDic": self.image_dict
        }
    
    def _generate_sketch_file(self, sketch_data, output_dir):
        """生成Sketch文件 - 使用独立的转换器"""
        from .json_to_sketch import JSONToSketchConverter
        
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        sketch_filename = f"converted_{uuid.uuid4().hex[:8]}.sketch"
        sketch_file_path = output_path / sketch_filename
        
        # 使用新的 JSON 到 Sketch 转换器
        converter = JSONToSketchConverter(verbose=self.verbose)
        success = converter.to_file(sketch_data, str(sketch_file_path))
        
        if success:
            self.log(f"Sketch 文件生成完成: {sketch_file_path}")
            return str(sketch_file_path)
        else:
            raise Exception("Sketch 文件生成失败")

def convert_ppt_to_sketch_async(task_id):
    """异步转换任务 - 使用增强版转换器"""
    from .models import ConversionTask
    from django.core.files import File
    
    task = None
    try:
        task = ConversionTask.objects.get(id=task_id)
        task.status = 'processing'
        task.save()
        
        logger.info(f"开始处理转换任务: {task_id}")
        
        # 检查文件是否存在
        if not os.path.exists(task.ppt_file.path):
            raise FileNotFoundError(f"PPT 文件不存在: {task.ppt_file.path}")
        
        # 执行转换
        converter = PPTToSketchConverter(verbose=True)
        output_dir = Path('media/outputs/sketch')
        
        sketch_file_path = converter.convert_ppt_to_sketch(
            task.ppt_file.path,
            output_dir
        )
        
        # 保存结果
        with open(sketch_file_path, 'rb') as f:
            task.sketch_file.save(
                os.path.basename(sketch_file_path),
                File(f),
                save=True
            )
        
        task.status = 'completed'
        task.error_message = None # 清除之前的错误信息
        task.save()
        
        # 清理临时文件
        if os.path.exists(sketch_file_path):
            os.remove(sketch_file_path)
        
        logger.info(f"转换任务完成: {task_id}")
        return True
        
    except Exception as e:
        if task:
            task.status = 'failed'
            task.error_message = str(e)
            task.save()
        logger.error(f"转换任务失败 {task_id}: {str(e)}", exc_info=True)
        return False