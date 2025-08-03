"""
增强的 PPT 转 Sketch 转换器
专门解决 Sketch 兼容性问题
"""

import os
import json
import zipfile
import uuid
import base64
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import logging

logger = logging.getLogger(__name__)

class EnhancedPPTToSketchConverter:
    """增强的 PPT 转 Sketch 转换器"""
    
    def __init__(self, artboard_width=1080, artboard_height=1920, verbose=False):
        """
        初始化增强转换器
        
        Args:
            artboard_width: 默认画板宽度（适合移动设备）
            artboard_height: 默认画板高度
            verbose: 详细日志
        """
        self.artboard_width = artboard_width
        self.artboard_height = artboard_height
        self.verbose = verbose
        self.image_dict = {}
        self.layer_counter = 0
        
    def log(self, message, level='info'):
        """日志记录"""
        if self.verbose:
            print(f"[{level.upper()}] {message}")
            
    def extract_color(self, color_obj):
        """提取颜色 - 增强版"""
        try:
            if isinstance(color_obj, RGBColor):
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": color_obj[2] / 255.0,
                    "green": color_obj[1] / 255.0,
                    "red": color_obj[0] / 255.0
                }
            elif hasattr(color_obj, 'rgb') and color_obj.rgb:
                rgb = color_obj.rgb
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": rgb[2] / 255.0,
                    "green": rgb[1] / 255.0,
                    "red": rgb[0] / 255.0
                }
            elif isinstance(color_obj, (tuple, list)) and len(color_obj) >= 3:
                return {
                    "_class": "color",
                    "alpha": 1,
                    "blue": color_obj[2] / 255.0,
                    "green": color_obj[1] / 255.0,
                    "red": color_obj[0] / 255.0
                }
        except Exception as e:
            self.log(f"颜色提取失败: {str(e)}", 'warning')
        
        # 默认黑色
        return {
            "_class": "color",
            "alpha": 1,
            "blue": 0,
            "green": 0,
            "red": 0
        }
    
    def optimize_image(self, image_blob, max_size=1024):
        """优化图片大小"""
        try:
            # 使用PIL压缩图片
            image = Image.open(io.BytesIO(image_blob))
            
            # 如果图片太大，进行压缩
            if max(image.size) > max_size:
                image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
                
            # 转换为RGB（如果是RGBA）
            if image.mode == 'RGBA':
                # 创建白色背景
                background = Image.new('RGB', image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[-1] if len(image.split()) == 4 else None)
                image = background
            elif image.mode != 'RGB':
                image = image.convert('RGB')
            
            # 保存为JPEG以减小文件大小
            output = io.BytesIO()
            image.save(output, format='JPEG', quality=85, optimize=True)
            optimized_data = output.getvalue()
            
            self.log(f"图片优化: {len(image_blob)} -> {len(optimized_data)} 字节")
            return optimized_data
            
        except Exception as e:
            self.log(f"图片优化失败: {str(e)}", 'warning')
            return image_blob
    
    def create_enhanced_image_layer(self, shape, layer_name):
        """创建增强的图片图层"""
        try:
            self.layer_counter += 1
            
            # 提取位置和尺寸
            x = shape.left / 914400 if hasattr(shape, 'left') else 0
            y = shape.top / 914400 if hasattr(shape, 'top') else 0
            width = shape.width / 914400 if hasattr(shape, 'width') else 100
            height = shape.height / 914400 if hasattr(shape, 'height') else 100
            
            # 缩放到合适的画板尺寸
            if width < 10:  # 太小的尺寸
                scale_factor = max(self.artboard_width / width / 10, self.artboard_height / height / 10)
                x *= scale_factor
                y *= scale_factor
                width *= scale_factor
                height *= scale_factor
            
            rotation = 0
            if hasattr(shape, 'rotation'):
                rotation = -shape.rotation
            
            # 处理图片数据
            image_id = str(uuid.uuid4())
            image_filename = f"{image_id}.jpg"
            
            if hasattr(shape, 'image'):
                try:
                    image_blob = shape.image.blob
                    # 优化图片
                    optimized_blob = self.optimize_image(image_blob)
                    # 转换为base64
                    image_b64 = base64.b64encode(optimized_blob).decode('utf-8')
                    self.image_dict[f"images/{image_filename}"] = image_b64
                    
                except Exception as img_e:
                    self.log(f"图片处理失败: {str(img_e)}", 'error')
                    return None
            
            # 创建Sketch图片图层
            layer = {
                "_class": "bitmap",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {
                    "_class": "rect",
                    "constrainProportions": True,
                    "height": height,
                    "width": width,
                    "x": x,
                    "y": y
                },
                "rotation": rotation,
                "isVisible": True,
                "isLocked": False,
                "hasClippingMask": False,
                "clippingMaskMode": 0,
                "userInfo": None,
                "image": {
                    "_class": "MSJSONFileReference",
                    "_ref_class": "MSImageData",
                    "_ref": f"images/{image_filename}"
                },
                "fillReplacesImage": False,
                "style": {
                    "_class": "style",
                    "endDecorationType": 0,
                    "miterLimit": 10,
                    "startDecorationType": 0,
                    "windingRule": 1
                }
            }
            
            self.log(f"创建增强图片图层: {layer_name}")
            return layer
            
        except Exception as e:
            self.log(f"创建图片图层失败: {str(e)}", 'error')
            return None
    
    def create_enhanced_group_layer(self, shape, layer_name):
        """创建增强的组图层"""
        try:
            self.layer_counter += 1
            
            # 提取位置和尺寸
            x = shape.left / 914400 if hasattr(shape, 'left') else 0
            y = shape.top / 914400 if hasattr(shape, 'top') else 0
            width = shape.width / 914400 if hasattr(shape, 'width') else 100
            height = shape.height / 914400 if hasattr(shape, 'height') else 100
            
            # 缩放处理
            if width < 10:
                scale_factor = max(self.artboard_width / width / 10, self.artboard_height / height / 10)
                x *= scale_factor
                y *= scale_factor
                width *= scale_factor
                height *= scale_factor
            
            rotation = 0
            if hasattr(shape, 'rotation'):
                rotation = -shape.rotation
            
            # 处理组内子形状
            sub_layers = []
            if hasattr(shape, 'shapes'):
                for i, sub_shape in enumerate(shape.shapes):
                    sub_layer = self.process_shape_enhanced(sub_shape, f"{layer_name}_item_{i}")
                    if sub_layer:
                        # 调整子图层相对位置
                        if 'frame' in sub_layer:
                            sub_layer['frame']['x'] -= x
                            sub_layer['frame']['y'] -= y
                        sub_layers.append(sub_layer)
            
            # 如果组为空，创建一个占位矩形
            if not sub_layers:
                placeholder = {
                    "_class": "rectangle",
                    "do_objectID": str(uuid.uuid4()),
                    "name": "Placeholder",
                    "frame": {
                        "_class": "rect",
                        "constrainProportions": False,
                        "height": height,
                        "width": width,
                        "x": 0,
                        "y": 0
                    },
                    "rotation": 0,
                    "isVisible": True,
                    "isLocked": False,
                    "style": {
                        "_class": "style",
                        "fills": [{
                            "_class": "fill",
                            "isEnabled": True,
                            "color": self.extract_color((200, 200, 200)),
                            "fillType": 0
                        }]
                    }
                }
                sub_layers.append(placeholder)
            
            # 创建组图层
            layer = {
                "_class": "group",
                "do_objectID": str(uuid.uuid4()),
                "name": layer_name,
                "frame": {
                    "_class": "rect",
                    "constrainProportions": False,
                    "height": height,
                    "width": width,
                    "x": x,
                    "y": y
                },
                "rotation": rotation,
                "isVisible": True,
                "isLocked": False,
                "hasClippingMask": False,
                "clippingMaskMode": 0,
                "userInfo": None,
                "layers": sub_layers,
                "style": {
                    "_class": "style",
                    "endDecorationType": 0,
                    "miterLimit": 10,
                    "startDecorationType": 0,
                    "windingRule": 1
                }
            }
            
            self.log(f"创建增强组图层: {layer_name}, {len(sub_layers)} 个子图层")
            return layer
            
        except Exception as e:
            self.log(f"创建组图层失败: {str(e)}", 'error')
            return None
    
    def process_shape_enhanced(self, shape, layer_name):
        """增强的形状处理"""
        try:
            shape_type = shape.shape_type
            
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self.create_enhanced_image_layer(shape, layer_name)
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                return self.create_enhanced_group_layer(shape, layer_name)
            else:
                # 其他类型暂时跳过
                self.log(f"跳过不支持的形状类型: {shape_type}", 'warning')
                return None
                
        except Exception as e:
            self.log(f"处理形状失败: {str(e)}", 'error')
            return None
    
    def convert_ppt_to_sketch_enhanced(self, ppt_path, output_dir):
        """增强的PPT转换"""
        try:
            self.log(f"开始增强转换: {ppt_path}")
            
            presentation = Presentation(ppt_path)
            
            # 获取PPT尺寸
            if hasattr(presentation, 'slide_width') and hasattr(presentation, 'slide_height'):
                original_width = presentation.slide_width / 914400
                original_height = presentation.slide_height / 914400
                
                # 如果原始尺寸太小，放大到合适的尺寸
                if original_width < 100 or original_height < 100:
                    scale_factor = max(self.artboard_width / original_width, self.artboard_height / original_height) * 0.8
                    self.artboard_width = original_width * scale_factor
                    self.artboard_height = original_height * scale_factor
                else:
                    self.artboard_width = original_width
                    self.artboard_height = original_height
                
                self.log(f"调整画板尺寸: {self.artboard_width:.1f} x {self.artboard_height:.1f}")
            
            # 转换所有幻灯片
            artboards = []
            for slide_index, slide in enumerate(presentation.slides):
                artboard = self.convert_slide_enhanced(slide, slide_index)
                if artboard:
                    artboards.append(artboard)
            
            if not artboards:
                raise Exception("没有成功转换任何幻灯片")
            
            # 创建Sketch数据结构
            sketch_data = self.create_enhanced_sketch_document(artboards)
            
            # 生成文件
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)
            
            sketch_filename = f"enhanced_{uuid.uuid4().hex[:8]}.sketch"
            sketch_file_path = output_path / sketch_filename
            
            # 使用JSONToSketchConverter
            from .json_to_sketch import JSONToSketchConverter
            converter = JSONToSketchConverter(verbose=self.verbose)
            success = converter.to_file(sketch_data, str(sketch_file_path))
            
            if success:
                self.log(f"增强转换完成: {sketch_file_path}")
                return str(sketch_file_path)
            else:
                raise Exception("文件生成失败")
                
        except Exception as e:
            self.log(f"增强转换失败: {str(e)}", 'error')
            raise e
    
    def convert_slide_enhanced(self, slide, slide_index):
        """增强的幻灯片转换"""
        try:
            artboard_name = f"Slide {slide_index + 1}"
            layers = []
            
            # 添加背景
            background_layer = {
                "_class": "rectangle",
                "do_objectID": str(uuid.uuid4()),
                "name": "Background",
                "frame": {
                    "_class": "rect",
                    "constrainProportions": False,
                    "height": self.artboard_height,
                    "width": self.artboard_width,
                    "x": 0,
                    "y": 0
                },
                "rotation": 0,
                "isVisible": True,
                "isLocked": False,
                "style": {
                    "_class": "style",
                    "fills": [{
                        "_class": "fill",
                        "isEnabled": True,
                        "color": self.extract_color((255, 255, 255)),
                        "fillType": 0
                    }]
                }
            }
            layers.append(background_layer)
            
            # 处理所有形状
            for shape_index, shape in enumerate(slide.shapes):
                layer_name = f"Element_{shape_index + 1}"
                layer = self.process_shape_enhanced(shape, layer_name)
                if layer:
                    layers.append(layer)
            
            # 创建画板
            artboard = {
                "_class": "artboard",
                "do_objectID": str(uuid.uuid4()),
                "name": artboard_name,
                "frame": {
                    "_class": "rect",
                    "constrainProportions": False,
                    "height": self.artboard_height,
                    "width": self.artboard_width,
                    "x": 0,
                    "y": slide_index * (self.artboard_height + 100)
                },
                "rotation": 0,
                "isVisible": True,
                "isLocked": False,
                "hasClippingMask": False,
                "layers": layers,
                "backgroundColor": self.extract_color((255, 255, 255)),
                "hasBackgroundColor": True,
                "includeInCloudUpload": True,
                "style": {
                    "_class": "style",
                    "endDecorationType": 0,
                    "miterLimit": 10,
                    "startDecorationType": 0,
                    "windingRule": 1
                }
            }
            
            self.log(f"转换幻灯片 {slide_index + 1}: {len(layers)} 个图层")
            return artboard
            
        except Exception as e:
            self.log(f"转换幻灯片失败: {str(e)}", 'error')
            return None
    
    def create_enhanced_sketch_document(self, artboards):
        """创建增强的Sketch文档"""
        
        pages_data = []
        
        for artboard in artboards:
            page_id = str(uuid.uuid4())
            page = {
                "_class": "page",
                "do_objectID": page_id,
                "name": artboard.get("name", "Page 1"),
                "layers": [artboard],
                "frame": {
                    "_class": "rect",
                    "constrainProportions": False,
                    "height": self.artboard_height,
                    "width": self.artboard_width,
                    "x": 0,
                    "y": 0
                },
                "style": {
                    "_class": "style"
                }
            }
            pages_data.append(page)
        
        document = {
            "_class": "document",
            "do_objectID": str(uuid.uuid4()),
            "pages": [],
            "layerStyles": {
                "_class": "sharedStyleContainer",
                "objects": []
            },
            "textStyles": {
                "_class": "sharedTextStyleContainer",
                "objects": []
            },
            "colors": {
                "_class": "swatchContainer",
                "objects": []
            },
            "gradients": {
                "_class": "gradientContainer",
                "objects": []
            },
            "colorSpace": 0,
            "currentPageIndex": 0
        }
        
        meta = {
            "_class": "meta",
            "app": "com.bohemiancoding.sketch3",
            "version": 134,
            "appVersion": "99.1",
            "build": 158899,
            "variant": "NONAPPSTORE",
            "created": {
                "_class": "MSJSONFileReference",
                "_ref_class": "MSDocumentData",
                "_ref": "document"
            },
            "pagesAndArtboards": {}
        }
        
        # 添加页面和画板信息
        for page in pages_data:
            page_id = page["do_objectID"]
            meta["pagesAndArtboards"][page_id] = {
                "name": page["name"],
                "artboards": {}
            }
            
            for artboard in page["layers"]:
                if artboard.get("_class") == "artboard":
                    artboard_id = artboard["do_objectID"]
                    meta["pagesAndArtboards"][page_id]["artboards"][artboard_id] = {
                        "name": artboard["name"]
                    }
        
        user = {
            "document": {
                "pageListHeight": 85,
                "pageListCollapsed": 0
            }
        }
        
        return {
            "contents": {
                "document": document,
                "pages": pages_data,
                "meta": meta,
                "user": user,
                "workspace": {}
            },
            "imageDic": self.image_dict
        } 